"""Integration helpers for openxml_audit.

Provides convenient utilities for integrating validation into various workflows:
- Context managers for validation
- Decorators for python-pptx integration
- pytest fixtures for testing
"""

from __future__ import annotations

import functools
import tempfile
from contextlib import contextmanager
from pathlib import Path
from typing import TYPE_CHECKING, Any, Callable, Generator, TypeVar

from openxml_audit.errors import FileFormat, ValidationResult
from openxml_audit.validator import OpenXmlValidator

if TYPE_CHECKING:
    pass

F = TypeVar("F", bound=Callable[..., Any])


@contextmanager
def validation_context(
    file_format: FileFormat = FileFormat.OFFICE_2019,
    max_errors: int = 100,
    raise_on_invalid: bool = False,
    strict: bool = True,
) -> Generator[OpenXmlValidator, None, None]:
    """Context manager for validation operations.

    Provides a configured validator instance for use within a context.

    Args:
        file_format: Office version to validate against.
        max_errors: Maximum number of errors to collect.
        raise_on_invalid: If True, raise ValueError on invalid files.

    Yields:
        Configured OpenXmlValidator instance.

    Example:
        from openxml_audit.helpers import validation_context

        with validation_context(raise_on_invalid=True) as validator:
            result = validator.validate("presentation.pptx")
            print(f"Valid: {result.is_valid}")
    """
    validator = OpenXmlValidator(
        file_format=file_format,
        max_errors=max_errors,
        strict=strict,
    )

    class ValidatorWrapper:
        """Wrapper that optionally raises on invalid files."""

        def __init__(self, inner: OpenXmlValidator) -> None:
            self._inner = inner

        def validate(self, path: str | Path) -> ValidationResult:
            result = self._inner.validate(path)
            if raise_on_invalid and not result.is_valid:
                error_summary = "; ".join(e.description for e in result.errors[:3])
                if len(result.errors) > 3:
                    error_summary += f"... (+{len(result.errors) - 3} more)"
                raise ValueError(f"Invalid PPTX: {error_summary}")
            return result

        def is_valid(self, path: str | Path) -> bool:
            return self._inner.is_valid(path)

        @property
        def file_format(self) -> FileFormat:
            return self._inner.file_format

        @property
        def max_errors(self) -> int:
            return self._inner.max_errors

    yield ValidatorWrapper(validator)  # type: ignore[misc]


def validate_on_save(
    file_format: FileFormat = FileFormat.OFFICE_2019,
    raise_on_invalid: bool = True,
) -> Callable[[F], F]:
    """Decorator to validate PPTX files created by python-pptx.

    Wraps functions that use python-pptx to create presentations,
    validating the output file after save.

    Args:
        file_format: Office version to validate against.
        raise_on_invalid: If True, raise ValueError on invalid output.

    Returns:
        Decorated function.

    Example:
        from pptx import Presentation
        from openxml_audit.helpers import validate_on_save

        @validate_on_save()
        def create_presentation(output_path: str) -> None:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            prs.save(output_path)

        # Will validate after save
        create_presentation("output.pptx")
    """

    def decorator(func: F) -> F:
        @functools.wraps(func)
        def wrapper(*args: Any, **kwargs: Any) -> Any:
            result = func(*args, **kwargs)

            # Try to find output path from args/kwargs
            output_path = None

            # Check common argument patterns
            if args and isinstance(args[-1], (str, Path)):
                candidate = Path(args[-1])
                if candidate.suffix.lower() == ".pptx":
                    output_path = candidate
            elif "output_path" in kwargs:
                output_path = Path(kwargs["output_path"])
            elif "path" in kwargs:
                output_path = Path(kwargs["path"])
            elif "filename" in kwargs:
                output_path = Path(kwargs["filename"])

            if output_path and output_path.exists():
                validator = OpenXmlValidator(file_format=file_format)
                validation_result = validator.validate(output_path)

                if not validation_result.is_valid and raise_on_invalid:
                    error_summary = "; ".join(
                        e.description for e in validation_result.errors[:3]
                    )
                    if len(validation_result.errors) > 3:
                        error_summary += f"... (+{len(validation_result.errors) - 3} more)"
                    raise ValueError(f"Generated invalid PPTX: {error_summary}")

            return result

        return wrapper  # type: ignore[return-value]

    return decorator


def require_valid_pptx(
    file_format: FileFormat = FileFormat.OFFICE_2019,
) -> Callable[[F], F]:
    """Decorator to validate PPTX input files before processing.

    Wraps functions that process PPTX files, validating the input
    file before the function runs.

    Args:
        file_format: Office version to validate against.

    Returns:
        Decorated function.

    Example:
        from openxml_audit.helpers import require_valid_pptx

        @require_valid_pptx()
        def process_presentation(input_path: str) -> dict:
            # Only runs if input is valid
            # ... process the file ...
            return {"slides": 10}

        # Will validate input first
        result = process_presentation("input.pptx")
    """

    def decorator(func: F) -> F:
        @functools.wraps(func)
        def wrapper(*args: Any, **kwargs: Any) -> Any:
            # Try to find input path from args/kwargs
            input_path = None

            if args and isinstance(args[0], (str, Path)):
                candidate = Path(args[0])
                if candidate.suffix.lower() == ".pptx":
                    input_path = candidate
            elif "input_path" in kwargs:
                input_path = Path(kwargs["input_path"])
            elif "path" in kwargs:
                input_path = Path(kwargs["path"])
            elif "filename" in kwargs:
                input_path = Path(kwargs["filename"])

            if input_path and input_path.exists():
                validator = OpenXmlValidator(file_format=file_format)
                validation_result = validator.validate(input_path)

                if not validation_result.is_valid:
                    error_summary = "; ".join(
                        e.description for e in validation_result.errors[:3]
                    )
                    if len(validation_result.errors) > 3:
                        error_summary += f"... (+{len(validation_result.errors) - 3} more)"
                    raise ValueError(f"Invalid input PPTX: {error_summary}")

            return func(*args, **kwargs)

        return wrapper  # type: ignore[return-value]

    return decorator


# pytest fixtures (import these in conftest.py)

def pytest_openxml_audit():
    """pytest fixture providing a validator instance.

    Usage in conftest.py:
        from openxml_audit.helpers import pytest_openxml_audit
        openxml_audit = pytest_openxml_audit

    Usage in tests:
        def test_my_pptx(openxml_audit):
            result = openxml_audit.validate("test.pptx")
            assert result.is_valid
    """
    import pytest

    @pytest.fixture
    def openxml_audit() -> OpenXmlValidator:
        """Fixture providing an OpenXmlValidator instance."""
        return OpenXmlValidator()

    return openxml_audit


def pytest_valid_pptx_path():
    """pytest fixture providing path to a minimal valid PPTX.

    Usage in conftest.py:
        from openxml_audit.helpers import pytest_valid_pptx_path
        valid_pptx_path = pytest_valid_pptx_path

    Usage in tests:
        def test_valid_file(valid_pptx_path):
            assert valid_pptx_path.exists()
    """
    import pytest

    @pytest.fixture
    def valid_pptx_path(tmp_path: Path) -> Generator[Path, None, None]:
        """Fixture providing a minimal valid PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        pptx_path = tmp_path / "valid.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(pptx_path))

        yield pptx_path

    return valid_pptx_path


def pytest_assert_valid_pptx():
    """pytest fixture providing assertion helper for PPTX validation.

    Usage in conftest.py:
        from openxml_audit.helpers import pytest_assert_valid_pptx
        assert_valid_pptx = pytest_assert_valid_pptx

    Usage in tests:
        def test_my_generator(assert_valid_pptx, tmp_path):
            output = tmp_path / "output.pptx"
            generate_pptx(output)
            assert_valid_pptx(output)
    """
    import pytest

    @pytest.fixture
    def assert_valid_pptx() -> Callable[[Path | str], None]:
        """Fixture providing validation assertion helper."""
        validator = OpenXmlValidator()

        def _assert(path: Path | str) -> None:
            result = validator.validate(path)
            if not result.is_valid:
                errors = "\n".join(f"  - {e.description}" for e in result.errors[:10])
                if len(result.errors) > 10:
                    errors += f"\n  ... (+{len(result.errors) - 10} more)"
                pytest.fail(f"PPTX validation failed:\n{errors}")

        return _assert

    return assert_valid_pptx
