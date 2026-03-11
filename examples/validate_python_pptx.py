#!/usr/bin/env python3
"""Generate a presentation with python-pptx and validate it.

Requirements:
    pip install openxml-audit python-pptx
"""

from pathlib import Path
from tempfile import TemporaryDirectory

from pptx import Presentation
from pptx.util import Inches, Pt

from openxml_audit import OpenXmlValidator, FileFormat


def generate_presentation(path: Path) -> None:
    """Create a sample presentation with python-pptx."""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Quarterly Report"
    slide.placeholders[1].text = "Generated with python-pptx"

    # Content slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Metrics"
    body = slide.placeholders[1]
    body.text = "Revenue: $1.2M"
    body.text_frame.paragraphs[0].font.size = Pt(18)

    # Slide with a shape
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))

    prs.save(str(path))


def main() -> None:
    with TemporaryDirectory() as tmp:
        path = Path(tmp) / "report.pptx"

        print("Generating presentation...")
        generate_presentation(path)

        print("Validating...")
        validator = OpenXmlValidator(file_format=FileFormat.OFFICE_2019)
        result = validator.validate(path)

        if result.is_valid:
            print(f"Valid! ({result.warning_count} warnings)")
        else:
            print(f"Invalid: {result.error_count} errors, {result.warning_count} warnings")
            for error in result.errors:
                print(f"  [{error.severity.value}] {error.description}")
                if error.part_uri:
                    print(f"    Part: {error.part_uri}")


if __name__ == "__main__":
    main()
